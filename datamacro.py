from pptx import Presentation # 라이브러리
import win32com.client # pptx 의 jpg 변환을 위한 라이브러리
import platedata # 번호판 정보를 담고 있는 라이브러리
import random
import pickle # 딕셔너리 저장용 라이브러리

# prs.slide_layouts[11] : 일반 번호판 레이아웃
# prs.slide_layouts[12] : 세종 번호판 레이아웃

n = 30 # 제작할 데이터셋의 수

Directory = 'D:\motorcycle\Plate\\'  # 생성된 pptx 파일의 저장 경로
Directoryjpg = 'D:\motorcycle\PlateJPG\\' # 생성된 jpg 파일의 저장 경로

def newNormalPlate(Id, dept_Office, dept_Local, id_char, id_num): # Id : 번호판의 순서 매기기용 숫자

    prs = Presentation('motorplate.pptx') # 파워포인트 객체 선언, 모체는 번호판 레이아웃 보유중인 motorplate.pptx

    normal_slide_layout = prs.slide_layouts[11] # 일반 번호판의 경우
    slide = prs.slides.add_slide(normal_slide_layout) 

    deptOffice = slide.placeholders[10] # 관할구청 [서울, 경기 , ...]
    deptOffice.text = dept_Office

    deptLocal = slide.placeholders[11] # 주소지 [전주시, 중구, ...] ※ 마지막글자는 포함하지 않는다!! (북구x 북o)
    deptLocal.text = dept_Local

    idchar = slide.placeholders[12] # 식별기호 [가~하]
    idchar.text = id_char

    idnum = slide.placeholders[13] # 식별숫자 [랜덤한 4자리 숫자]
    id_num = str(id_num)
    idnum.text = id_num

    plate_Id = str(Id) + '_' # n_ 꼴로 변환

    filename = plate_Id + dept_Office + dept_Local + id_char + id_num # 출력된 파일명은 1_전북전주가7777.pptx 의 꼴.
    prs.save(Directory + filename + '.pptx')
    return(filename)
    


def newSejongPlate(Id, id_char, id_num):

    prs = Presentation('motorplate.pptx') # 파워포인트 객체 선언, 모체는 번호판 레이아웃 보유중인 motorplate.pptx


    sej_slide_layout = prs.slide_layouts[12] # 세종 번호판의 경우에는 관할구청이 "세종" 으로, 주소지는 없다.
    slide = prs.slides.add_slide(sej_slide_layout)

    idchar = slide.placeholders[12] # 식별기호 [가~하]
    idchar.text = id_char

    idnum = slide.placeholders[13] # 식별숫자 [랜덤한 4자리 숫자]
    id_num = str(id_num)
    idnum.text = id_num

    plate_Id = str(Id) + '_' # n_ 꼴로 변환

    filename = plate_Id + '세종' + id_char + id_num # 출력된 파일명은 1_세종가7777.pptx 의 꼴.
    prs.save(Directory + filename + '.pptx')
    return(filename)

filenames = [] # 파일명 저장 리스트
filedir = [] # ppt 파일경로 저장 리스트
exportdir = [] # jpg 파일경로 저장 리스트

for i in range(n): # pptx 제작하기

    office = random.choice(platedata.office)

    if office != '세종':

        local = random.choice(platedata.local)
        local = local[:-1]
        gatoha = random.choice(platedata.gatoha)
        num = random.randrange(1000,9999)

        filenames.append(newNormalPlate(i,office,local,gatoha,num))
        filedir.append(Directory + filenames[i] +'.pptx')
        exportdir.append(Directoryjpg + filenames[i]+'.jpg')

    else:

        gatoha = random.choice(platedata.gatoha)
        num = random.randrange(1000,9999)

        filenames.append(newSejongPlate(i,gatoha,num))
        filedir.append(Directory + filenames[i]+'.pptx')
        exportdir.append(Directoryjpg + newSejongPlate(i,gatoha,num)+'.jpg')


for i in range(len(filedir)): #제작한 pptx 를 jpg 로 변환
    Application = win32com.client.Dispatch("PowerPoint.Application")
    Presentation = Application.Presentations.Open(filedir[i])
    Presentation.Slides[0].Export(exportdir[i], "JPG")
    Application.Quit()
    Presentation =  None
    Application = None

label = []

for i in range(len(filenames)):

    filename = filenames[i]

    info = {}

    image = {}
    image['file_name'] = filename+'.jpg'

    while filename[0] != '_': # 파일명 숫자부분 제거
        filename = filename[1:]


    temp = []
    temp.append(filename[1:3]) # 관할구청 추가
    filename = filename[3:]
    temp.append(filename[-4:]) # 번호 추가
    filename = filename[:-4]
    temp.insert(1, filename[-1]) 
    filename = filename[:-1]
    
    if  temp: 
        temp.insert(1, filename)

    text = {}
    text['output'] = temp

    info['image'] = image
    info['text'] = text

    label.append(info)

print(label)

with open('label.pkl','wb') as f: # 생성된 라벨 저장용
    pickle.dump(label, f)

